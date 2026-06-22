from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


OVERVIEW_OUTPUT_FILE = Path(__file__).with_name("agentic-approach-overview.pptx")

NAVY = RGBColor(15, 23, 42)
SLATE = RGBColor(30, 41, 59)
MUTED = RGBColor(100, 116, 139)
TEAL = RGBColor(13, 148, 136)
TEAL_DARK = RGBColor(15, 118, 110)
AMBER = RGBColor(217, 119, 6)
ROSE = RGBColor(190, 24, 93)
CREAM = RGBColor(248, 250, 252)
WHITE = RGBColor(255, 255, 255)
PALE_TEAL = RGBColor(204, 251, 241)
PALE_AMBER = RGBColor(254, 243, 199)
PALE_ROSE = RGBColor(251, 207, 232)
PALE_SLATE = RGBColor(226, 232, 240)

FONT_HEAD = "Aptos Display"
FONT_BODY = "Aptos"
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)
BRAND_LABEL = ".claude agentic workflow"


def add_background(slide, color):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        0,
        0,
        SLIDE_WIDTH,
        SLIDE_HEIGHT,
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    slide.shapes._spTree.remove(shape._element)
    slide.shapes._spTree.insert(2, shape._element)


def add_textbox(
    slide,
    left,
    top,
    width,
    height,
    text,
    *,
    font_size=18,
    font_name=FONT_BODY,
    color=SLATE,
    bold=False,
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
    margin=0.08,
):
    textbox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    textbox.text_frame.clear()
    textbox.text_frame.word_wrap = True
    textbox.text_frame.margin_left = Inches(margin)
    textbox.text_frame.margin_right = Inches(margin)
    textbox.text_frame.margin_top = Inches(margin)
    textbox.text_frame.margin_bottom = Inches(margin)
    textbox.text_frame.vertical_anchor = valign
    paragraph = textbox.text_frame.paragraphs[0]
    paragraph.alignment = align
    run = paragraph.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return textbox


def add_paragraphs(
    slide,
    left,
    top,
    width,
    height,
    lines,
    *,
    font_size=18,
    color=SLATE,
    line_spacing=1.15,
):
    textbox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    frame = textbox.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(0.08)
    frame.margin_right = Inches(0.08)
    frame.margin_top = Inches(0.08)
    frame.margin_bottom = Inches(0.08)
    for index, line in enumerate(lines):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.alignment = PP_ALIGN.LEFT
        paragraph.line_spacing = line_spacing
        run = paragraph.add_run()
        run.text = line
        run.font.name = FONT_BODY
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
    return textbox


def add_card(slide, left, top, width, height, fill_color, title, body_lines, *, title_color=SLATE, body_color=SLATE):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = fill_color
    add_textbox(
        slide,
        left + 0.18,
        top + 0.14,
        width - 0.36,
        0.42,
        title,
        font_size=20,
        font_name=FONT_HEAD,
        color=title_color,
        bold=True,
    )
    add_paragraphs(
        slide,
        left + 0.18,
        top + 0.62,
        width - 0.36,
        height - 0.78,
        body_lines,
        font_size=15,
        color=body_color,
        line_spacing=1.12,
    )


def add_chip(slide, left, top, width, text, fill_color, text_color):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(0.4),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = fill_color
    add_textbox(
        slide,
        left,
        top + 0.02,
        width,
        0.32,
        text,
        font_size=12,
        color=text_color,
        bold=True,
        align=PP_ALIGN.CENTER,
    )


def add_footer(slide, text, dark=False):
    add_textbox(
        slide,
        0.45,
        7.02,
        12.4,
        0.28,
        text,
        font_size=10,
        color=WHITE if dark else MUTED,
    )


def add_brand_badge(slide, dark=False):
    fill_color = WHITE if dark else NAVY
    text_color = NAVY if dark else WHITE
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(10.55),
        Inches(0.42),
        Inches(2.15),
        Inches(0.38),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = fill_color
    add_textbox(
        slide,
        10.55,
        0.44,
        2.15,
        0.28,
        BRAND_LABEL,
        font_size=11,
        color=text_color,
        bold=True,
        align=PP_ALIGN.CENTER,
    )


def add_speaker_notes(slide, title, bullets):
    frame = slide.notes_slide.notes_text_frame
    frame.clear()
    frame.word_wrap = True
    paragraph = frame.paragraphs[0]
    paragraph.text = title
    for bullet in bullets:
        note_paragraph = frame.add_paragraph()
        note_paragraph.text = bullet


def set_table_cell(cell, text, *, color=SLATE, size=14, bold=False, fill=None):
    if fill is not None:
        cell.fill.solid()
        cell.fill.fore_color.rgb = fill
    frame = cell.text_frame
    frame.clear()
    frame.word_wrap = True
    paragraph = frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.LEFT
    run = paragraph.add_run()
    run.text = text
    run.font.name = FONT_BODY
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def build_title_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, NAVY)
    add_brand_badge(slide, dark=True)
    accent = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(0), Inches(0), Inches(13.333), Inches(0.22)
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = TEAL
    accent.line.fill.background()
    add_chip(slide, 0.6, 0.55, 2.35, "AGENTIC DELIVERY MODEL", PALE_TEAL, TEAL_DARK)
    add_textbox(
        slide,
        0.62,
        1.25,
        8.7,
        1.4,
        "How this .claude workflow operates and where the value actually comes from",
        font_size=28,
        font_name=FONT_HEAD,
        color=WHITE,
        bold=True,
    )
    add_textbox(
        slide,
        0.65,
        2.65,
        7.8,
        1.0,
        "This deck is based on the current agent set, specialist skills, and durable handoff artifacts in the workspace as of 2026-06-20.",
        font_size=18,
        color=RGBColor(226, 232, 240),
    )
    add_card(
        slide,
        0.65,
        4.35,
        3.0,
        1.55,
        WHITE,
        "What it is",
        [
            "A routed multi-agent workflow with specialist roles instead of one broad prompt.",
        ],
        title_color=NAVY,
    )
    add_card(
        slide,
        3.88,
        4.35,
        3.0,
        1.55,
        PALE_TEAL,
        "How it works",
        [
            "Intake, routing, optional design, targeted execution, then independent validation.",
        ],
        title_color=TEAL_DARK,
    )
    add_card(
        slide,
        7.11,
        4.35,
        2.9,
        1.55,
        PALE_AMBER,
        "Why it matters",
        [
            "It reduces wasted work, routing ambiguity, and context loss on non-trivial tasks.",
        ],
        title_color=AMBER,
    )
    add_card(
        slide,
        10.23,
        4.35,
        2.45,
        1.55,
        PALE_ROSE,
        "What to remember",
        [
            "The value is controlled specialization with durable evidence, not agent count by itself.",
        ],
        title_color=ROSE,
    )
    add_footer(slide, "Source surfaces: agents/task-router.md, agents/feature-intake-reviewer.md, skills/agent-handoff-evidence-best-practices/SKILL.md", dark=True)
    add_speaker_notes(
        slide,
        "Open by framing this as a governed delivery model, not a novelty multi-agent demo.",
        [
            "Explain that the deck is grounded in the current repo configuration rather than generic AI workflow theory.",
            "Set the expectation that the real value comes from control points, specialist boundaries, and durable artifacts.",
            "Use the four summary cards to preview the talk: what it is, how it works, why it matters, and what to remember.",
        ],
    )


def build_problem_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, CREAM)
    add_brand_badge(slide)
    add_chip(slide, 0.55, 0.45, 1.7, "WHY THIS EXISTS", PALE_SLATE, NAVY)
    add_textbox(slide, 0.55, 0.95, 7.6, 0.7, "The model is solving concrete execution failure modes", font_size=26, font_name=FONT_HEAD, color=NAVY, bold=True)
    add_textbox(slide, 0.58, 1.62, 8.4, 0.5, "The repo is not using agents just to fan out work. It is adding gates and durable state where single-threaded chat usually breaks down.", font_size=16, color=MUTED)
    add_card(
        slide,
        0.65,
        2.15,
        5.9,
        4.35,
        WHITE,
        "Failure modes without structure",
        [
            "- vague requests go straight into implementation",
            "- one agent makes design, coding, testing, and review calls at the same time",
            "- context is trapped in chat history instead of preserved as artifacts",
            "- quality checks become optional when time is tight",
            "- cross-stack work creates routing churn before real work even starts",
        ],
        title_color=NAVY,
    )
    add_card(
        slide,
        6.78,
        2.15,
        5.9,
        4.35,
        PALE_TEAL,
        "Repo design choices that address them",
        [
            "- feature-intake-reviewer adds a scope, cost-band, and approval gate",
            "- task-router chooses the shortest valid specialist chain",
            "- architect and schema agents are only introduced when uncertainty is real",
            "- testers and reviewers are separate specialist roles",
            "- handoff and evidence files preserve decisions and validation outside the chat",
        ],
        title_color=TEAL_DARK,
    )
    add_footer(slide, "The operating model trades a small amount of upfront structure for better control on multi-step work.")
    add_speaker_notes(
        slide,
        "Describe the specific failure modes this workflow is designed to reduce.",
        [
            "Emphasize that the system is solving for ambiguity, context loss, and inconsistent validation.",
            "Point out that the fix is not more agents by itself; it is putting the right gates before and after implementation.",
            "Use this slide to justify why a single long prompt is often the wrong control model for non-trivial work.",
        ],
    )


def build_workflow_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, NAVY)
    add_brand_badge(slide, dark=True)
    add_chip(slide, 0.58, 0.48, 1.55, "HOW IT FLOWS", PALE_TEAL, TEAL_DARK)
    add_textbox(slide, 0.6, 0.98, 8.4, 0.7, "The workflow uses targeted stages, not a single long-running agent", font_size=26, font_name=FONT_HEAD, color=WHITE, bold=True)
    steps = [
        ("1. Intake", "Clarify scope, risk, cost band, and whether approval is required.", TEAL),
        ("2. Route", "Pick the minimum specialist sequence that can complete the task safely.", RGBColor(14, 165, 233)),
        ("3. Design if needed", "Use architecture or schema specialists only when the shape of the change is unresolved.", AMBER),
        ("4. Build", "Hand implementation to the backend, frontend, refactor, or wiki specialist.", RGBColor(99, 102, 241)),
        ("5. Validate", "Run dedicated testers and read-only reviewers as a quality gate.", ROSE),
        ("6. Persist evidence", "Capture routing, handoff, review, and test results so the next agent does not start cold.", RGBColor(34, 197, 94)),
    ]
    left = 0.6
    for title, body, color in steps:
        add_card(slide, left, 2.1, 2.0, 3.5, WHITE, title, [body], title_color=color)
        left += 2.08
    add_textbox(slide, 0.62, 6.05, 12.1, 0.55, "The key principle is conditional depth: the workflow expands only when the task actually needs more planning, schema work, testing, or review.", font_size=17, color=RGBColor(226, 232, 240))
    add_footer(slide, "This is why the router explicitly keeps the chain as short as possible while still honoring intake, planning, schema, and validation needs.", dark=True)
    add_speaker_notes(
        slide,
        "Walk the audience across the stages as a control flow, not a rigid assembly line.",
        [
            "Call out that intake and routing are there to prevent downstream waste, not to slow work down.",
            "Stress the conditional design step: architecture and schema specialists only appear when the task has unresolved design risk.",
            "Close by noting that validation and persistent evidence are first-class stages rather than optional cleanup at the end.",
        ],
    )


def build_agent_map_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, CREAM)
    add_brand_badge(slide)
    add_chip(slide, 0.55, 0.45, 1.6, "AGENT SYSTEM", PALE_SLATE, NAVY)
    add_textbox(slide, 0.55, 0.95, 7.5, 0.7, "Current role map: 13 agents with explicit boundaries", font_size=26, font_name=FONT_HEAD, color=NAVY, bold=True)
    rows = [
        (0.68, "Orchestration", PALE_TEAL, TEAL_DARK, ["feature-intake-reviewer", "task-router"]),
        (1.82, "Design and data", PALE_AMBER, AMBER, ["react-dotnet-solution-architect", "dotnet-efcore-schema-designer", "dotnet-efcore-migrations"]),
        (2.96, "Execution", WHITE, NAVY, ["dotnet-backend-developer", "react-frontend-developer", "react-dotnet-refactor-specialist"]),
        (4.10, "Quality", PALE_ROSE, ROSE, ["dotnet-backend-unit-tester", "react-frontend-unit-tester", "dotnet-backend-code-reviewer", "react-frontend-code-reviewer"]),
        (5.24, "Knowledge", PALE_SLATE, SLATE, ["llm-wiki-builder"]),
    ]
    for top, title, fill, title_color, agents in rows:
        shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.62), Inches(top), Inches(12.0), Inches(0.88))
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
        shape.line.color.rgb = fill
        add_textbox(slide, 0.82, top + 0.16, 2.05, 0.3, title, font_size=18, font_name=FONT_HEAD, color=title_color, bold=True)
        add_textbox(slide, 3.0, top + 0.15, 9.15, 0.38, " | ".join(agents), font_size=16, color=SLATE if fill != PALE_TEAL else TEAL_DARK)
    add_textbox(slide, 0.66, 6.45, 12.0, 0.45, "The naming scheme matters. Stack-specific prefixes reduce routing ambiguity and make it obvious when a task belongs to React, .NET, EF Core, refactor, quality, or wiki work.", font_size=15, color=MUTED)
    add_footer(slide, "Reviewers are intentionally read-only. Testers are distinct from implementation roles. Those boundaries are part of the value, not incidental details.")
    add_speaker_notes(
        slide,
        "Use this slide to explain why explicit role boundaries improve execution quality.",
        [
            "Point out that orchestration, design, implementation, quality, and knowledge management are separated on purpose.",
            "Mention that read-only reviewers and dedicated testers reduce self-review bias and create cleaner validation gates.",
            "Highlight that stack-aware naming lowers the routing cost before any implementation work begins.",
        ],
    )


def build_artifacts_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, CREAM)
    add_brand_badge(slide)
    add_chip(slide, 0.55, 0.45, 2.0, "MEMORY AND GOVERNANCE", PALE_SLATE, NAVY)
    add_textbox(slide, 0.55, 0.95, 7.8, 0.7, "Durable artifacts keep the workflow from depending on chat memory", font_size=26, font_name=FONT_HEAD, color=NAVY, bold=True)
    add_card(
        slide,
        0.65,
        1.95,
        6.0,
        4.95,
        WHITE,
        "Artifact layer",
        [
            "- feature-intake.md freezes scope, assumptions, cost band, and approval state",
            "- task-routing.md records the ordered agent chain and its gates",
            "- architecture-plan.md preserves decisions that should not be rediscovered later",
            "- task-handoff.md gives the next agent a current execution snapshot",
            "- review-findings.md and test-report.md make quality evidence durable",
        ],
        title_color=NAVY,
    )
    add_card(
        slide,
        6.78,
        1.95,
        5.9,
        4.95,
        PALE_TEAL,
        "Skill layer",
        [
            "- 16 skill directories centralize standards instead of duplicating them in every agent",
            "- examples include React component guidance, state management, TypeScript standards, .NET architecture, EF Core migration safety, and Web API security",
            "- the result is repeatability: specialists can stay focused while still applying shared best practices",
        ],
        title_color=TEAL_DARK,
    )
    add_footer(slide, "Artifacts preserve task state. Skills preserve domain standards. Together they keep handoffs cheap and consistent.")
    add_speaker_notes(
        slide,
        "Explain that the repo has both a memory layer and a standards layer.",
        [
            "The markdown artifacts hold task state so the next agent does not need to reconstruct intent from chat history.",
            "The skill directories hold reusable best practices so specialists can stay focused without losing consistency.",
            "This combination is what makes the workflow resumable, auditable, and repeatable across sessions.",
        ],
    )


def build_example_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, NAVY)
    add_brand_badge(slide, dark=True)
    add_chip(slide, 0.58, 0.48, 1.85, "EXAMPLE PATH", PALE_TEAL, TEAL_DARK)
    add_textbox(slide, 0.6, 0.98, 8.2, 0.7, "How a cross-stack feature moves through the system", font_size=26, font_name=FONT_HEAD, color=WHITE, bold=True)
    add_textbox(slide, 0.62, 1.58, 10.2, 0.42, "Example request: Add a new export flow across the API and UI, but only if the scope is small enough to justify the work.", font_size=16, color=RGBColor(226, 232, 240))
    cards = [
        (0.65, 2.2, 2.0, 1.6, "Intake", ["Clarify scope.", "Estimate cost band.", "Ask for approval if needed."], PALE_TEAL, TEAL_DARK),
        (2.8, 2.2, 2.0, 1.6, "Routing", ["Choose the smallest chain.", "Skip architecture if the change is obvious."], WHITE, NAVY),
        (4.95, 2.2, 2.0, 1.6, "Optional design", ["Use architect or schema roles only when contracts or data are unsettled."], PALE_AMBER, AMBER),
        (7.1, 2.2, 2.0, 1.6, "Implementation", ["Backend and frontend specialists build the approved slice."], WHITE, NAVY),
        (9.25, 2.2, 1.95, 1.6, "Quality", ["Dedicated test and review steps verify the result."], PALE_ROSE, ROSE),
        (11.35, 2.2, 1.35, 1.6, "Artifacts", ["Persist handoff and evidence."], PALE_SLATE, SLATE),
    ]
    for left, top, width, height, title, lines, fill, title_color in cards:
        add_card(slide, left, top, width, height, fill, title, lines, title_color=title_color)
    add_card(
        slide,
        0.7,
        4.55,
        3.9,
        1.4,
        WHITE,
        "Value created at the front",
        ["The intake gate prevents an expensive chain from starting on a fuzzy request."],
        title_color=TEAL_DARK,
    )
    add_card(
        slide,
        4.72,
        4.55,
        3.9,
        1.4,
        PALE_AMBER,
        "Value created in the middle",
        ["Conditional routing adds design depth only when the change genuinely needs it."],
        title_color=AMBER,
    )
    add_card(
        slide,
        8.74,
        4.55,
        3.9,
        1.4,
        PALE_ROSE,
        "Value created at the end",
        ["Testing, review, and artifacts make the result easier to trust, resume, and audit."],
        title_color=ROSE,
    )
    add_footer(slide, "This is the practical meaning of the model: early gating, conditional specialization, and durable evidence.", dark=True)
    add_speaker_notes(
        slide,
        "Use the export-flow example to show how the system spends coordination only where it pays back.",
        [
            "Start with the intake gate and explain that the request can still be stopped, clarified, or split before expensive execution starts.",
            "Show that architecture is optional rather than automatic, which keeps the workflow lean when the change is already well understood.",
            "Finish on the end-state: separate quality checks plus durable artifacts make the output easier to trust and resume.",
        ],
    )


def build_value_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, CREAM)
    add_brand_badge(slide)
    add_chip(slide, 0.55, 0.45, 1.9, "ACTUAL VALUE", PALE_SLATE, NAVY)
    add_textbox(slide, 0.55, 0.95, 7.6, 0.7, "Where the value really comes from", font_size=26, font_name=FONT_HEAD, color=NAVY, bold=True)
    positions = [
        (0.65, 1.95, PALE_TEAL, "Less wasted spend", ["Intake and routing prevent long execution chains from starting on unclear requests."], TEAL_DARK),
        (4.37, 1.95, WHITE, "Better specialist fit", ["Agent boundaries map to real concerns: frontend, backend, schema, refactor, quality, and wiki work."], NAVY),
        (8.09, 1.95, PALE_AMBER, "Clearer handoffs", ["Artifacts replace brittle chat-only memory, so the next agent can continue without reconstruction."], AMBER),
        (0.65, 4.15, WHITE, "Stronger quality gates", ["Testers and read-only reviewers are distinct roles, which reduces self-review bias."], NAVY),
        (4.37, 4.15, PALE_ROSE, "Reusable expertise", ["Skills centralize standards once and let multiple specialists apply them consistently."], ROSE),
        (8.09, 4.15, PALE_SLATE, "Lower routing ambiguity", ["Stack-aware naming makes it obvious whether work belongs to React, .NET, EF Core, refactor, or wiki lanes."], SLATE),
    ]
    for left, top, fill, title, lines, title_color in positions:
        add_card(slide, left, top, 3.45, 1.7, fill, title, lines, title_color=title_color)
    add_footer(slide, "Important nuance: this value shows up on non-trivial work. For truly tiny one-step tasks, a long chain would just add overhead.")
    add_speaker_notes(
        slide,
        "Anchor the value discussion in concrete operational outcomes.",
        [
            "Less wasted spend comes from stopping fuzzy work before it becomes a long multi-agent sequence.",
            "Clearer handoffs and reusable expertise reduce rediscovery and make specialist work more predictable.",
            "Be explicit that this value is most visible on cross-stack or high-confidence work rather than tiny one-off edits.",
        ],
    )


def build_tradeoff_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, CREAM)
    add_brand_badge(slide)
    add_chip(slide, 0.55, 0.45, 2.15, "ECONOMICS AND TRADEOFFS", PALE_SLATE, NAVY)
    add_textbox(slide, 0.55, 0.95, 8.1, 0.7, "Why this is better than one broad agent on complex work", font_size=26, font_name=FONT_HEAD, color=NAVY, bold=True)
    table = slide.shapes.add_table(7, 3, Inches(0.65), Inches(1.85), Inches(12.0), Inches(4.8)).table
    table.columns[0].width = Inches(2.5)
    table.columns[1].width = Inches(4.55)
    table.columns[2].width = Inches(4.95)
    headers = ["Decision area", "Single general agent", "This agentic model"]
    for index, text in enumerate(headers):
        set_table_cell(table.cell(0, index), text, color=WHITE, size=15, bold=True, fill=NAVY)
    rows = [
        ("Best fit", "Trivial one-off edits or quick answers.", "Multi-step work where routing, validation, or persistence matter."),
        ("Upfront overhead", "Very low.", "Higher, but deliberate and usually limited to intake and routing."),
        ("Handling ambiguity", "Often resolves ambiguity by making assumptions in-flight.", "Pulls ambiguity forward through intake, planning, and explicit approval."),
        ("Specialization depth", "Generic unless the prompt is very detailed.", "Explicitly aligned to frontend, backend, schema, refactor, review, test, and wiki concerns."),
        ("Quality control", "Validation is easy to skip or merge into implementation.", "Dedicated tester and reviewer roles make quality a formal gate."),
        ("Resumability", "Context stays in chat and is easy to lose.", "Artifacts make the task resumable and auditable across agents and sessions."),
    ]
    for row_index, row in enumerate(rows, start=1):
        fill = WHITE if row_index % 2 else PALE_SLATE
        for column_index, text in enumerate(row):
            set_table_cell(table.cell(row_index, column_index), text, fill=fill, bold=column_index == 0)
    add_footer(slide, "The design goal is not to maximize agents. It is to spend coordination only where coordination pays back.")
    add_speaker_notes(
        slide,
        "Use this comparison to prevent over-selling the model.",
        [
            "Acknowledge that a single general agent is still the right tool for very small requests or quick answers.",
            "Position this workflow as an operating model for ambiguity, validation, and resumability rather than a universal default.",
            "That honesty makes the ROI argument stronger because it is tied to task shape instead of hype.",
        ],
    )


def build_closing_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, NAVY)
    add_brand_badge(slide, dark=True)
    add_chip(slide, 0.58, 0.48, 1.65, "BOTTOM LINE", PALE_TEAL, TEAL_DARK)
    add_textbox(slide, 0.6, 1.0, 8.1, 0.8, "The value is controlled specialization with durable evidence", font_size=28, font_name=FONT_HEAD, color=WHITE, bold=True)
    add_paragraphs(
        slide,
        0.62,
        1.95,
        6.2,
        2.2,
        [
            "- use intake when scope, cost, or approval is uncertain",
            "- route non-trivial work instead of relying on an implicit all-purpose prompt",
            "- keep implementation, testing, and review as distinct concerns",
            "- treat the markdown artifacts as the workflow memory layer",
        ],
        font_size=18,
        color=RGBColor(226, 232, 240),
        line_spacing=1.18,
    )
    add_card(
        slide,
        7.15,
        1.95,
        5.05,
        1.45,
        WHITE,
        "What success looks like",
        ["Fewer vague starts, fewer routing detours, cleaner handoffs, and stronger trust in what shipped."],
        title_color=NAVY,
    )
    add_card(
        slide,
        7.15,
        3.68,
        5.05,
        1.45,
        PALE_TEAL,
        "Where to be disciplined",
        ["Do not force a long agent chain onto tiny tasks. The workflow should expand only when the task complexity justifies it."],
        title_color=TEAL_DARK,
    )
    add_card(
        slide,
        7.15,
        5.41,
        5.05,
        1.15,
        PALE_AMBER,
        "One sentence summary",
        ["This approach turns AI work from a long chat into a governed delivery system."],
        title_color=AMBER,
    )
    add_footer(slide, "Generated from the current .claude agent and skill definitions on 2026-06-20.", dark=True)
    add_speaker_notes(
        slide,
        "Close by summarizing the model in plain business language.",
        [
            "The workflow turns AI execution from a long chat into a governed delivery system with clear stages and durable evidence.",
            "Reinforce the discipline point: do not force a large chain onto tiny work that does not need it.",
            "If the audience is evaluating adoption, suggest starting with cross-stack or high-risk tasks where the control benefits are easiest to measure.",
        ],
    )


def main():
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT
    build_title_slide(prs)
    build_problem_slide(prs)
    build_workflow_slide(prs)
    build_agent_map_slide(prs)
    build_artifacts_slide(prs)
    build_example_slide(prs)
    build_value_slide(prs)
    build_tradeoff_slide(prs)
    build_closing_slide(prs)
    prs.save(OVERVIEW_OUTPUT_FILE)
    print(f"Created {OVERVIEW_OUTPUT_FILE.name} with {len(prs.slides)} slides")


if __name__ == "__main__":
    main()